#!/usr/bin/env python3
"""
File Parser: Convert various file formats to plain text
Supports: TXT, PDF, DOCX, CSV, JSON, HTML, MD, XLS/XLSX, PPT/PPTX
"""

import sys
import json
import os
from pathlib import Path

def parse_txt(file_path):
    """Parse plain text file"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        return f"Error reading TXT file: {str(e)}"

def parse_json(file_path):
    """Parse JSON file to readable text"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        return f"Error parsing JSON: {str(e)}"

def parse_csv(file_path):
    """Parse CSV file"""
    try:
        import csv
        rows = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.DictReader(f)
            for row in reader:
                rows.append(row)
        
        if not rows:
            return "CSV file is empty"
        
        # Format as readable text
        text = "CSV Data:\n"
        text += "=" * 80 + "\n"
        for idx, row in enumerate(rows, 1):
            text += f"Row {idx}:\n"
            for key, value in row.items():
                text += f"  {key}: {value}\n"
            text += "\n"
        return text
    except Exception as e:
        return f"Error parsing CSV: {str(e)}"

def parse_html(file_path):
    """Parse HTML file - extract text content"""
    try:
        from html.parser import HTMLParser
        
        class HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
                self.in_script = False
                self.in_style = False
            
            def handle_starttag(self, tag, attrs):
                if tag.lower() in ('script', 'style'):
                    self.in_script = tag.lower() == 'script'
                    self.in_style = tag.lower() == 'style'
            
            def handle_endtag(self, tag):
                if tag.lower() in ('script', 'style'):
                    self.in_script = False
                    self.in_style = False
                elif tag.lower() in ('p', 'div', 'br', 'li'):
                    self.text.append('\n')
            
            def handle_data(self, data):
                if not self.in_script and not self.in_style:
                    text = data.strip()
                    if text:
                        self.text.append(text)
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        
        extractor = HTMLTextExtractor()
        extractor.feed(html_content)
        return '\n'.join(extractor.text)
    except Exception as e:
        return f"Error parsing HTML: {str(e)}"

def parse_pdf(file_path):
    """Parse PDF file - requires PyPDF2"""
    try:
        import PyPDF2
        text = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(f"\n--- Page {page_num} ---\n")
                        text.append(page_text)
                except Exception as e:
                    text.append(f"\n--- Page {page_num} (error: {str(e)}) ---\n")
        return ''.join(text) if text else "No text found in PDF"
    except ImportError:
        return "PyPDF2 not installed. Install with: pip install PyPDF2"
    except Exception as e:
        return f"Error parsing PDF: {str(e)}"

def parse_docx(file_path):
    """Parse DOCX file - requires python-docx"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        return '\n'.join(text) if text else "No text found in DOCX"
    except ImportError:
        return "python-docx not installed. Install with: pip install python-docx"
    except Exception as e:
        return f"Error parsing DOCX: {str(e)}"

def parse_markdown(file_path):
    """Parse Markdown file"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        return f"Error reading Markdown: {str(e)}"

def parse_excel(file_path):
    """Parse Excel file - requires openpyxl or pandas"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path)
        text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text.append(f"\n=== Sheet: {sheet_name} ===\n")
            for row in ws.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    text.append(row_text)
                    text.append('\n')
        return ''.join(text) if text else "No data found in Excel"
    except ImportError:
        return "openpyxl not installed. Install with: pip install openpyxl"
    except Exception as e:
        return f"Error parsing Excel: {str(e)}"

def parse_powerpoint(file_path):
    """Parse PowerPoint file - requires python-pptx"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text.append(f"\n=== Slide {slide_num} ===\n")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text.append(shape.text)
                    text.append('\n')
        return ''.join(text) if text else "No text found in PowerPoint"
    except ImportError:
        return "python-pptx not installed. Install with: pip install python-pptx"
    except Exception as e:
        return f"Error parsing PowerPoint: {str(e)}"

def parse_file(file_path):
    """Route file parsing based on extension"""
    if not os.path.exists(file_path):
        return json.dumps({"error": f"File not found: {file_path}"})
    
    file_ext = Path(file_path).suffix.lower()
    
    parsers = {
        '.txt': parse_txt,
        '.md': parse_markdown,
        '.json': parse_json,
        '.csv': parse_csv,
        '.html': parse_html,
        '.htm': parse_html,
        '.pdf': parse_pdf,
        '.docx': parse_docx,
        '.xlsx': parse_excel,
        '.xls': parse_excel,
        '.pptx': parse_powerpoint,
        '.ppt': parse_powerpoint,
    }
    
    parser = parsers.get(file_ext, parse_txt)
    
    try:
        content = parser(file_path)
        result = {
            "success": True,
            "filename": os.path.basename(file_path),
            "file_type": file_ext,
            "content": content,
            "char_count": len(content),
            "token_estimate": len(content) // 4  # Rough estimation
        }
    except Exception as e:
        result = {
            "success": False,
            "error": str(e)
        }
    
    return json.dumps(result, ensure_ascii=False)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "No file path provided"}))
        sys.exit(1)
    
    file_path = sys.argv[1]
    print(parse_file(file_path))
